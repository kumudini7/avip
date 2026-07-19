from __future__ import annotations

import math
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Spacer,
    Paragraph,
    Table,
    TableStyle,
)
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import settings
from app.models.analysis import ProcessAnalysis
from app.models.document import ProcessDocument
from app.models.project import Project
from app.models.proposal import Proposal
from app.services.analysis_engine import analyze_process
from app.services.document_processor import normalize_text
from app.services.document_processor import ensure_directory

ACTION_KEYWORDS = {
    "rpa": {"download", "upload", "copy", "paste", "update", "reconcile", "transfer", "route", "export", "import", "enter", "post"},
    "ai": {"extract", "classify", "summarize", "analyze", "read", "detect", "understand", "interpret", "predict"},
    "human": {"approve", "review", "verify", "judgment", "escalate", "sign off", "signoff", "exception"},
}

INDUSTRY_KPIS = {
    "finance": ["Invoice Processing Time", "Payment Delay", "Accuracy", "DSO"],
    "hr": ["Hiring Time", "Employee Satisfaction", "Payroll Accuracy"],
    "sales": ["Lead Response Time", "Conversion Rate", "Customer Retention"],
    "operations": ["Cycle Time", "Throughput", "Rework"],
    "healthcare": ["Case Turnaround Time", "Compliance", "Accuracy"],
    "manufacturing": ["Cycle Time", "Throughput", "Quality"],
}

BENCHMARKS = {
    "healthcare": {"automation": 68, "roi": 240, "payback": "8 months"},
    "manufacturing": {"automation": 82, "roi": 310, "payback": "6 months"},
    "finance": {"automation": 74, "roi": 285, "payback": "7 months"},
    "hr": {"automation": 64, "roi": 210, "payback": "9 months"},
}


def _keywords_present(text: str, keywords: set[str]) -> int:
    lower = text.lower()
    return sum(1 for keyword in keywords if keyword in lower)


def _split_steps(text: str) -> list[str]:
    normalized = normalize_text(text)
    if not normalized:
        return []

    raw_segments = re.split(r"[\n.;]+", normalized)
    steps: list[str] = []
    for segment in raw_segments:
        clean = re.sub(r"^[\-\*\d\)\(]+", "", segment).strip()
        if len(clean) < 4:
            continue
        if any(word in clean.lower() for word in ("step", "process", "receive", "extract", "approve", "validate", "update", "review", "download", "upload", "email")):
            steps.append(clean)

    if not steps:
        sentences = re.split(r"(?<=[.!?])\s+", normalized)
        steps = [sentence.strip() for sentence in sentences if len(sentence.strip()) > 4]

    return steps[:12]


def _classify_step(step: str) -> dict[str, Any]:
    lower = step.lower()
    rpa_hits = _keywords_present(lower, ACTION_KEYWORDS["rpa"])
    ai_hits = _keywords_present(lower, ACTION_KEYWORDS["ai"])
    human_hits = _keywords_present(lower, ACTION_KEYWORDS["human"])

    if human_hits and (ai_hits or rpa_hits):
        recommendation = "Human + AI"
        rationale = "Requires review or judgment with automated support."
    elif ai_hits and rpa_hits:
        recommendation = "AI + RPA"
        rationale = "Combines document understanding with a transaction step."
    elif ai_hits:
        recommendation = "AI"
        rationale = "Language or document interpretation is the main effort."
    elif rpa_hits:
        recommendation = "RPA"
        rationale = "Repeatable rule-based transaction is a good automation fit."
    elif human_hits:
        recommendation = "Human Only"
        rationale = "The step appears to need direct human judgment."
    else:
        recommendation = "Workflow"
        rationale = "Best handled through orchestration and exception routing."

    confidence = min(97, 58 + (ai_hits + rpa_hits + human_hits) * 8)
    return {
        "step": step,
        "recommendation": recommendation,
        "rationale": rationale,
        "confidence": confidence,
    }


def _detect_industry(project: Project) -> str:
    candidate = " ".join(filter(None, [project.industry, project.name, project.description])).lower()
    for industry in INDUSTRY_KPIS:
        if industry in candidate:
            return industry
    return "operations"


def _score_analysis(text: str, documents: list[ProcessDocument], project: Project) -> dict[str, Any]:
    normalized = normalize_text(text)
    steps = _split_steps(normalized)
    classifications = [_classify_step(step) for step in steps]

    decision_points = len(re.findall(r"\b(if|when|otherwise|threshold|exception|above|below|approve|reject)\b", normalized.lower()))
    exceptions = len(re.findall(r"\b(exception|escalat|error|failure|manual review)\b", normalized.lower()))
    system_interactions = len(re.findall(r"\b(api|erp|sap|crm|email|portal|database|system)\b", normalized.lower()))
    document_signals = len(re.findall(r"\b(pdf|doc|docx|invoice|form|table|attachment|scan|image)\b", normalized.lower()))
    human_judgment = len([item for item in classifications if item["recommendation"] in {"Human Only", "Human + AI"}])
    ai_candidates = len([item for item in classifications if "AI" in item["recommendation"]])
    rpa_candidates = len([item for item in classifications if "RPA" in item["recommendation"]])

    complexity_parts = {
        "rules_complexity": min(20, decision_points * 4),
        "app_count": min(15, max(1, system_interactions) * 3),
        "decision_complexity": min(15, decision_points * 3),
        "document_complexity": min(15, document_signals * 2 + len(documents) * 2),
        "exception_rate": min(10, exceptions * 3),
        "ocr_requirement": 10 if any(doc.storage_path.lower().endswith((".pdf", ".png", ".jpg", ".jpeg")) for doc in documents) else 4,
        "nlp_requirement": 8 if ai_candidates else 3,
        "human_judgment": min(10, human_judgment * 3),
        "api_availability": 8 if system_interactions else 4,
    }
    complexity_score = min(100, sum(complexity_parts.values()))

    automation_bonus = min(12, rpa_candidates * 2 + ai_candidates * 2)
    readiness = max(35, min(97, 100 - complexity_score + automation_bonus))

    confidence = max(
        65,
        min(97, 55 + len(steps) * 3 + len(documents) * 5 + min(10, len(normalized) // 300)),
    )
    if not normalized:
        confidence = 45

    if readiness >= 85:
        risk_level = "Low"
    elif readiness >= 70:
        risk_level = "Medium"
    else:
        risk_level = "High"

    hours_saved = max(60, len(steps) * 42 + rpa_candidates * 28 + ai_candidates * 24)
    value_metrics = {
        "hours_saved": hours_saved,
        "quality_improvement": min(40, 12 + ai_candidates * 4 + document_signals * 2),
        "compliance_improvement": min(45, 15 + human_judgment * 5 + exceptions * 2),
        "customer_satisfaction": min(30, 10 + readiness // 4),
        "accuracy_improvement": min(40, 15 + ai_candidates * 5),
        "risk_reduction": min(35, 10 + max(0, 20 - complexity_score // 2)),
        "revenue_protection": f"INR {max(6, len(steps) * 2)} Lakhs",
    }

    industry = _detect_industry(project)
    kpis = INDUSTRY_KPIS.get(industry, INDUSTRY_KPIS["operations"])
    benchmark = BENCHMARKS.get(industry, {"automation": 70, "roi": 250, "payback": "8 months"})

    recommendations = {
        "automation_stack": [
            "RPA",
            "OCR",
            "LLM",
            "Workflow",
            "Human Approval",
        ],
        "vendor_stack": [
            "UiPath",
            "Azure OpenAI / OpenAI / Gemini",
            "Document Processing",
            "Local File Storage",
        ],
        "human_review_policy": "Review only high-value exceptions and low-confidence cases.",
    }

    scenario_data = [
        {"scenario": "Only RPA", "roi": max(120, readiness + 25), "payback": "12 months"},
        {"scenario": "RPA + OCR", "roi": max(165, readiness + 70), "payback": "8 months"},
        {"scenario": "RPA + LLM", "roi": max(220, readiness + 120), "payback": "5 months"},
    ]

    risks = [
        "Operational risk from manual dependency",
        "Compliance risk if approvals are bypassed",
        "Security risk for document handling",
        "Model hallucination risk on ambiguous inputs",
        "Business continuity risk if systems are unavailable",
    ]
    mitigations = [
        "Human approval for high-value cases",
        "Audit logs and traceability",
        "Confidence thresholds for AI outputs",
        "Secure local file storage with access control",
    ]

    return {
        "process_name": project.name,
        "industry": industry,
        "automation_readiness": readiness,
        "ai_confidence": confidence,
        "risk_level": risk_level,
        "summary": (
            f"{project.name} shows {readiness}% automation readiness with {confidence}% AI confidence. "
            f"The process appears to be a {risk_level.lower()}-risk candidate with strong potential for RPA, OCR, and LLM support."
        ),
        "steps": classifications,
        "complexity": {
            "total": complexity_score,
            "parts": complexity_parts,
        },
        "value_metrics": value_metrics,
        "recommendations": recommendations,
        "risks": risks,
        "mitigations": mitigations,
        "benchmark": benchmark,
        "scenarios": scenario_data,
        "kpi_mapping": kpis,
        "generated_at": datetime.utcnow().isoformat(),
    }


def build_analysis_payload(project: Project, documents: list[ProcessDocument]) -> dict[str, Any]:
    return analyze_process(project, documents)


def _proposal_directory(project_id: int) -> Path:
    return ensure_directory(Path(settings.report_dir) / f"project_{project_id}")


def _format_lines(lines: list[str]) -> list[str]:
    return [line for line in lines if line]


def build_proposal_payload(project: Project, analysis_payload: dict[str, Any]) -> dict[str, Any]:
    executive_summary = analysis_payload["summary"]
    proposal_payload = {
        "executive_summary": executive_summary,
        "automation_opportunities": analysis_payload["steps"],
        "architecture": analysis_payload["recommendations"],
        "roi": analysis_payload["scenarios"],
        "kpi_mapping": analysis_payload["kpi_mapping"],
        "risks": analysis_payload["risks"],
        "mitigations": analysis_payload["mitigations"],
        "benchmark": analysis_payload["benchmark"],
    }

    output_dir = _proposal_directory(project.id)
    pdf_path = output_dir / f"{project.name.lower().replace(' ', '_')}_proposal.pdf"
    ppt_path = output_dir / f"{project.name.lower().replace(' ', '_')}_proposal.pptx"

    _write_proposal_pdf(pdf_path, project, analysis_payload, proposal_payload)
    _write_proposal_pptx(ppt_path, project, analysis_payload, proposal_payload)

    proposal_payload["files"] = {
        "pdf": str(pdf_path),
        "pptx": str(ppt_path),
    }
    return proposal_payload


def _write_proposal_pdf(output_path: Path, project: Project, analysis_payload: dict[str, Any], proposal_payload: dict[str, Any]) -> None:
    doc = SimpleDocTemplate(str(output_path), pagesize=A4, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="SectionTitle", parent=styles["Heading2"], textColor=colors.HexColor("#0f766e")))
    styles.add(ParagraphStyle(name="Body", parent=styles["BodyText"], leading=14, fontSize=10))

    story: list[Any] = []
    story.append(Paragraph(f"{project.name} Proposal", styles["Title"]))
    story.append(Spacer(1, 0.18 * inch))
    story.append(Paragraph(proposal_payload["executive_summary"], styles["Body"]))
    story.append(Spacer(1, 0.2 * inch))

    metric_rows = [
        ["Automation Readiness", f'{analysis_payload["automation_readiness"]}%'],
        ["AI Confidence", f'{analysis_payload["ai_confidence"]}%'],
        ["Risk Level", analysis_payload["risk_level"]],
        ["Benchmark ROI", f'{analysis_payload["benchmark"]["roi"]}%'],
        ["Payback", analysis_payload["benchmark"]["payback"]],
    ]
    story.append(Table(metric_rows, colWidths=[2.3 * inch, 3.7 * inch], style=TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#d1fae5")),
        ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#0f172a")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#94a3b8")),
        ("BACKGROUND", (0, 1), (-1, -1), colors.white),
        ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("LEADING", (0, 0), (-1, -1), 12),
    ])))
    story.append(Spacer(1, 0.2 * inch))

    sections = [
        ("Automation Opportunities", [f'{item["step"]} - {item["recommendation"]}' for item in analysis_payload["steps"]]),
        ("Recommended Stack", proposal_payload["architecture"]["automation_stack"]),
        ("KPIs", proposal_payload["kpi_mapping"]),
        ("Risks", proposal_payload["risks"]),
        ("Mitigations", proposal_payload["mitigations"]),
    ]
    for title, items in sections:
        story.append(Paragraph(title, styles["SectionTitle"]))
        story.append(Spacer(1, 0.08 * inch))
        for item in _format_lines(items):
            story.append(Paragraph(f"- {item}", styles["Body"]))
        story.append(Spacer(1, 0.12 * inch))

    doc.build(story)


def _write_proposal_pptx(output_path: Path, project: Project, analysis_payload: dict[str, Any], proposal_payload: dict[str, Any]) -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{project.name} Proposal"
    title_slide.placeholders[1].text = analysis_payload["summary"]

    def add_bullets_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = ""
            paragraph.level = 0
            run = paragraph.add_run()
            run.text = bullet
            run.font.size = Pt(20)

    add_bullets_slide(
        "Automation Opportunities",
        [f'{item["step"]} - {item["recommendation"]}' for item in analysis_payload["steps"]],
    )
    add_bullets_slide(
        "Recommended Stack",
        proposal_payload["architecture"]["automation_stack"] + proposal_payload["architecture"]["vendor_stack"],
    )
    add_bullets_slide(
        "Risks and Mitigations",
        [*proposal_payload["risks"], *proposal_payload["mitigations"]],
    )
    add_bullets_slide(
        "Business Value",
        [f'{key.replace("_", " ").title()}: {value}' for key, value in analysis_payload["value_metrics"].items()],
    )

    prs.save(str(output_path))


def create_proposal(project: Project, analysis: ProcessAnalysis) -> tuple[dict[str, Any], str, str]:
    proposal_payload = build_proposal_payload(project, analysis.analysis_json)
    pdf_path = proposal_payload["files"]["pdf"]
    ppt_path = proposal_payload["files"]["pptx"]
    return proposal_payload, pdf_path, ppt_path
